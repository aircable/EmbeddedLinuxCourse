"""
Theme Engine for Embedded Linux Course Slides
Usage: from theme_engine import *
Run with: /opt/hermes/.venv/bin/python
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Theme Colors ─────────────────────────────────────────────────────────────
DARK_BG        = RGBColor(0x15, 0x23, 0x34)
LIGHT_BG       = RGBColor(0xF6, 0xF8, 0xFB)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT      = RGBColor(0x15, 0x23, 0x34)
GRAY_TEXT      = RGBColor(0x99, 0x99, 0x99)
MUTED_TEXT     = RGBColor(0x59, 0x59, 0x59)
FOOTER_BAR     = RGBColor(0xE8, 0xE8, 0xE8)
CODE_BG        = RGBColor(0x1E, 0x1E, 0x1E)
CODE_TEXT      = RGBColor(0xD4, 0xD4, 0xD4)
LAB_ACCENT     = RGBColor(0x2E, 0x7D, 0x32)
LAB_BG = RGBColor(0xE8, 0xF5, 0xE9)
WARN_BORDER = RGBColor(0xFF, 0x8F, 0x00)
WARN_BG = RGBColor(0xFF, 0xF3, 0xE0)

GOLD           = RGBColor(0xF0, 0xB3, 0x42)
ORANGE         = RGBColor(0xE6, 0x86, 0x2E)
TEAL           = RGBColor(0x19, 0x7A, 0x8C)
BLUE           = RGBColor(0x3D, 0x6F, 0xF6)
ORANGE_TINT    = RGBColor(0xFD, 0xF4, 0xE9)
TEAL_TINT      = RGBColor(0xE6, 0xF5, 0xF7)
BLUE_TINT      = RGBColor(0xEB, 0xF1, 0xFF)
GOLD_TINT      = RGBColor(0xFF, 0xF9, 0xEB)

ACCENTS   = [GOLD, ORANGE, TEAL, BLUE]
TINTS     = [ORANGE_TINT, TEAL_TINT, BLUE_TINT, GOLD_TINT]

# Legacy color aliases
MUTED = MUTED_TEXT
LAB_GREEN = LAB_ACCENT

# ── Fonts ────────────────────────────────────────────────────────────────────
FONT_HEADING  = "Arial"
FONT_BODY     = "Arial"
FONT_CODE     = "Courier New"

# Legacy/compatibility constant aliases used by older rebuild scripts
FONT_H = FONT_HEADING
FONT_B = FONT_BODY
FONT_C = FONT_CODE

# ── Sizes ────────────────────────────────────────────────────────────────────
SZ_TITLE    = Pt(36)
SZ_SUBTITLE = Pt(24)
SZ_HEADING  = Pt(20)
SZ_BODY     = Pt(14)
SZ_CODE     = Pt(11)
SZ_CAPTION  = Pt(10)
SZ_FOOTER   = Pt(9)

# Small text alias
SZ_SM = SZ_CAPTION

# ── Layout ──────────────────────────────────────────────────────────────────
SLIDE_W      = Inches(13.333)
SLIDE_H      = Inches(7.5)
MARGIN       = Inches(0.5)
ACCENT_H     = Inches(0.15)
FOOTER_H     = Inches(0.25)
CONTENT_TOP  = Inches(0.9)
CARD_PAD     = Inches(0.15)

# Legacy layout aliases
TITLE_TOP = Inches(0.3)
TITLE_H = Inches(0.7)
BODY_TOP = CONTENT_TOP

# ── Global state ─────────────────────────────────────────────────────────────
_acc_idx = 0

def accent(i=None):
    global _acc_idx
    if i is None:
        c = ACCENTS[_acc_idx % len(ACCENTS)]
        _acc_idx += 1
        return c
    return ACCENTS[i % len(ACCENTS)]

def tint(i=None):
    if i is None:
        return TINTS[_acc_idx % len(TINTS) - 1] if _acc_idx > 0 else TINTS[0]
    return TINTS[i % len(TINTS)]

def reset_accents():
    global _acc_idx
    _acc_idx = 0

# ── Slide Background ──────────────────────────────────────────────────────────
def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

# ── Accent Bar ───────────────────────────────────────────────────────────────
def accent_bar(slide, prs_or_color, color=None):
    if hasattr(prs_or_color, "slide_width"):
        prs = prs_or_color
        accent_color = color
    else:
        accent_color = prs_or_color
        prs = color
    if prs is None:
        raise ValueError("Presentation object is required for accent_bar")
    c = accent_color if accent_color else accent()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, ACCENT_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = c
    bar.line.fill.background()
    sp = bar._element
    parent = sp.getparent()
    parent.remove(sp)
    parent.insert(0, sp)

# ── Footer Bar ───────────────────────────────────────────────────────────────
def footer_bar(slide, prs, text=""):
    y = prs.slide_height - FOOTER_H
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, FOOTER_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = FOOTER_BAR
    bar.line.fill.background()
    if text:
        tf = bar.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = SZ_FOOTER
        p.font.color.rgb = MUTED_TEXT
        p.font.name = FONT_BODY

# ── Text Box ─────────────────────────────────────────────────────────────────
def textbox(slide, l, t, w, h, text, font=FONT_BODY, size=SZ_BODY,
            color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = font
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tb

# Compatibility wrappers for older slide scripts
def txt(slide, l, t, w, h, text, fn=FONT_BODY, sz=SZ_BODY, clr=DARK_TEXT,
        bold=False, align=PP_ALIGN.LEFT, wrap=True):
    return textbox(slide, l, t, w, h, text, font=fn, size=sz,
                   color=clr, bold=bold, align=align, wrap=wrap)

def footer(slide, prs, text=""):
    return footer_bar(slide, prs, text)

def add_card(slide, l, t, w, h, title, items, accent_c, tint_c):
    return card(slide, l, t, w, h, title, items, accent_c, tint_c)

# ── Bullet List ──────────────────────────────────────────────────────────────
def bullets(slide, l, t, w, h, items, font=FONT_BODY, size=SZ_BODY,
            color=DARK_TEXT, spacing=Pt(6)):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.name = font
        p.font.size = size
        p.font.color.rgb = color
        p.space_after = spacing
    return tb

# ── Code Block ───────────────────────────────────────────────────────────────
def code_block(slide, l, t, w, h, code):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = CODE_BG
    bg.line.fill.background()
    tb = slide.shapes.add_textbox(l + CARD_PAD, t + CARD_PAD,
                                   w - 2 * CARD_PAD, h - 2 * CARD_PAD)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.name = FONT_CODE
    p.font.size = SZ_CODE
    p.font.color.rgb = CODE_TEXT
    return tb

# ── Card ─────────────────────────────────────────────────────────────────────
def card(slide, l, t, w, h, title, items, accent_c, tint_c):
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    bg.fill.solid()
    bg.fill.fore_color.rgb = tint_c
    bg.line.fill.background()
    th = Inches(0.4)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, th)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_c
    bar.line.fill.background()
    tb = slide.shapes.add_textbox(l + CARD_PAD, t + Inches(0.05),
                                   w - 2 * CARD_PAD, th)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEADING
    p.font.size = SZ_HEADING
    p.font.color.rgb = WHITE
    p.font.bold = True
    ct = t + th + Inches(0.1)
    ch = h - th - Inches(0.15)
    bullets(slide, l + CARD_PAD, ct, w - 2 * CARD_PAD, ch, items,
            font=FONT_BODY, size=Pt(12))

# ── Title Slide ──────────────────────────────────────────────────────────────
def make_title_slide(prs, title, subtitle="", footer_text=""):
    reset_accents()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK_BG)
    accent_bar(slide, prs, GOLD)
    textbox(slide, MARGIN, Inches(1.5),
             SLIDE_W - 2 * MARGIN, Inches(1.5),
             title, font=FONT_HEADING, size=SZ_TITLE,
             color=WHITE, bold=True)
    if subtitle:
        textbox(slide, MARGIN, Inches(3.2),
                 SLIDE_W - 2 * MARGIN, Inches(1),
                 subtitle, font=FONT_HEADING, size=SZ_SUBTITLE,
                 color=GOLD, bold=True)
    if footer_text:
        textbox(slide, MARGIN, prs.slide_height - FOOTER_H - Inches(0.1),
                 SLIDE_W - 2 * MARGIN, Inches(0.3),
                 footer_text, size=Pt(11), color=GRAY_TEXT)
    footer_bar(slide, prs)
    return slide

# ── Section Header Slide ────────────────────────────────────────────────────
def make_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, DARK_BG)
    color = accent()
    accent_bar(slide, prs, color)
    textbox(slide, MARGIN, Inches(2),
             SLIDE_W - 2 * MARGIN, Inches(1.2),
             title, font=FONT_HEADING, size=Pt(40),
             color=WHITE, bold=True)
    if subtitle:
        textbox(slide, MARGIN, Inches(3.3),
                 SLIDE_W - 2 * MARGIN, Inches(0.8),
                 subtitle, font=FONT_BODY, size=SZ_SUBTITLE,
                 color=color)
    footer_bar(slide, prs)
    return slide

# ── Content Slide ────────────────────────────────────────────────────────────
def make_content_slide(prs, title, body_items=None, body_text=None,
                       accent_i=None, footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    accent_bar(slide, prs, accent(accent_i))
    textbox(slide, MARGIN, Inches(0.3),
             SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, font=FONT_HEADING, size=Pt(28),
             color=DARK_TEXT, bold=True)
    top = CONTENT_TOP
    h = prs.slide_height - top - FOOTER_H - Inches(0.2)
    if body_items and isinstance(body_items, list) and body_items:
        if isinstance(body_items[0], str):
            bullets(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, h, body_items)
        elif isinstance(body_items[0], tuple):
            n = len(body_items)
            card_h = h / n - Inches(0.1)
            card_w = (SLIDE_W - 2 * MARGIN - Inches(0.2)) / 2
            for idx, (heading, items) in enumerate(body_items):
                col = idx % 2
                row = idx // 2
                cl = MARGIN + col * (card_w + Inches(0.2))
                ct = top + row * (card_h + Inches(0.1))
                card(slide, cl, ct, card_w, card_h, heading, items,
                     accent(idx), tint(idx))
    elif body_text:
        textbox(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, h,
                body_text, font=FONT_BODY, size=SZ_BODY, color=DARK_TEXT)
    footer_bar(slide, prs, footer_text)
    return slide

# ── Two-Column Slide ─────────────────────────────────────────────────────────
def make_two_col_slide(prs, title, left_title, left_items,
                        right_title, right_items, footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    accent_bar(slide, prs)
    textbox(slide, MARGIN, Inches(0.3),
             SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, font=FONT_HEADING, size=Pt(28),
             color=DARK_TEXT, bold=True)
    col_w = (SLIDE_W - 2 * MARGIN - Inches(0.3)) / 2
    top = CONTENT_TOP
    h = prs.slide_height - top - FOOTER_H - Inches(0.2)
    textbox(slide, MARGIN, top, col_w, Inches(0.5),
            left_title, font=FONT_HEADING, size=SZ_HEADING,
            color=DARK_TEXT, bold=True)
    bullets(slide, MARGIN, top + Inches(0.5), col_w, h - Inches(0.5),
            left_items)
    l2 = MARGIN + col_w + Inches(0.3)
    textbox(slide, l2, top, col_w, Inches(0.5),
            right_title, font=FONT_HEADING, size=SZ_HEADING,
            color=DARK_TEXT, bold=True)
    bullets(slide, l2, top + Inches(0.5), col_w, h - Inches(0.5),
            right_items)
    footer_bar(slide, prs, footer_text)
    return slide

# ── Code Slide ───────────────────────────────────────────────────────────────
def make_code_slide(prs, title, code, accent_i=None, footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    accent_bar(slide, prs, accent(accent_i))
    textbox(slide, MARGIN, Inches(0.3),
             SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, font=FONT_HEADING, size=Pt(28),
             color=DARK_TEXT, bold=True)
    top = CONTENT_TOP
    h = prs.slide_height - top - FOOTER_H - Inches(0.2)
    code_block(slide, MARGIN, top, SLIDE_W - 2 * MARGIN, h, code)
    footer_bar(slide, prs, footer_text)
    return slide

# ── Lab Slide ────────────────────────────────────────────────────────────────
def make_lab_slide(prs, title, steps, accent_i=None, footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    accent_bar(slide, prs, LAB_ACCENT)
    textbox(slide, MARGIN, Inches(0.3),
             SLIDE_W - 2 * MARGIN, Inches(0.5),
             title, font=FONT_HEADING, size=Pt(28),
             color=LAB_ACCENT, bold=True)
    top = CONTENT_TOP
    h = prs.slide_height - top - FOOTER_H - Inches(0.2)
    step_h = h / len(steps)
    for i, (heading, detail) in enumerate(steps):
        y = top + i * step_h
        num = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, MARGIN, y + Inches(0.05),
            Inches(0.4), Inches(0.4))
        num.fill.solid()
        num.fill.fore_color.rgb = LAB_ACCENT
        num.line.fill.background()
        ntb = slide.shapes.add_textbox(
            MARGIN, y + Inches(0.08), Inches(0.4), Inches(0.35))
        ntf = ntb.text_frame
        np = ntf.paragraphs[0]
        np.text = str(i + 1)
        np.font.name = FONT_HEADING
        np.font.size = SZ_HEADING
        np.font.color.rgb = WHITE
        np.font.bold = True
        np.alignment = PP_ALIGN.CENTER
        textbox(slide, MARGIN + Inches(0.6), y,
                 SLIDE_W - 2 * MARGIN - Inches(0.6), step_h,
                 heading, font=FONT_HEADING, size=Pt(16),
                 color=DARK_TEXT, bold=True)
        textbox(slide, MARGIN + Inches(0.6), y + Inches(0.3),
                 SLIDE_W - 2 * MARGIN - Inches(0.6), step_h - Inches(0.3),
                 detail, font=FONT_BODY, size=Pt(13), color=DARK_TEXT)
    footer_bar(slide, prs, footer_text)
    return slide

# ── Flow Diagram Slide ───────────────────────────────────────────────────────
def make_flow_slide(prs, title, steps, accent_i=None, footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    accent_bar(slide, prs, accent(accent_i))
    textbox(slide, MARGIN, Inches(0.3),
             SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, font=FONT_HEADING, size=Pt(28),
             color=DARK_TEXT, bold=True)
    top = CONTENT_TOP
    h = prs.slide_height - top - FOOTER_H - Inches(0.3)
    box_h = h / len(steps) - Inches(0.08)
    box_w = SLIDE_W * 0.7
    box_l = (SLIDE_W - box_w) / 2
    for i, (name, desc) in enumerate(steps):
        y = top + i * (box_h + Inches(0.08))
        color = ACCENTS[i % len(ACCENTS)]
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, box_l, y, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        textbox(slide, box_l + Inches(0.15), y + Inches(0.05),
                 box_w - Inches(0.3), box_h * 0.5,
                 name, font=FONT_HEADING, size=Pt(18),
                 color=WHITE, bold=True)
        textbox(slide, box_l + Inches(0.15), y + box_h * 0.45,
                 box_w - Inches(0.3), box_h * 0.5,
                 desc, font=FONT_BODY, size=Pt(12), color=WHITE)
        if i < len(steps) - 1:
            ax = SLIDE_W / 2 - Inches(0.15)
            ay = y + box_h
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.ISOSCELES_TRIANGLE, ax, ay,
                Inches(0.3), Inches(0.08))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = GRAY_TEXT
            arrow.line.fill.background()
    footer_bar(slide, prs, footer_text)
    return slide

# ── Summary Slide ────────────────────────────────────────────────────────────
def make_summary_slide(prs, title, items, accent_i=None, footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    accent_bar(slide, prs, accent(accent_i))
    textbox(slide, MARGIN, Inches(0.3),
             SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, font=FONT_HEADING, size=Pt(28),
             color=DARK_TEXT, bold=True)
    top = CONTENT_TOP
    h = prs.slide_height - top - FOOTER_H - Inches(0.2)
    item_h = h / len(items)
    for i, item in enumerate(items):
        y = top + i * item_h
        color = ACCENTS[i % len(ACCENTS)]
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, MARGIN, y,
            Inches(0.08), item_h - Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        textbox(slide, MARGIN + Inches(0.2), y,
                 SLIDE_W - 2 * MARGIN - Inches(0.25), item_h,
                 item, font=FONT_BODY, size=Pt(16), color=DARK_TEXT)
    footer_bar(slide, prs, footer_text)
    return slide

# ── Comparison Table Slide ───────────────────────────────────────────────────
def make_comparison_slide(prs, title, headers, rows, accent_i=None, footer_text=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg(slide, LIGHT_BG)
    accent_bar(slide, prs, accent(accent_i))
    textbox(slide, MARGIN, Inches(0.3),
             SLIDE_W - 2 * MARGIN, Inches(0.7),
             title, font=FONT_HEADING, size=Pt(28),
             color=DARK_TEXT, bold=True)
    top = CONTENT_TOP
    h = prs.slide_height - top - FOOTER_H - Inches(0.2)
    ncols = len(headers)
    col_w = (SLIDE_W - 2 * MARGIN) / ncols
    for c_idx, header in enumerate(headers):
        l = MARGIN + c_idx * col_w
        cell = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            l, top, col_w - Inches(0.05), Inches(0.4))
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG
        cell.line.fill.background()
        tb = slide.shapes.add_textbox(
            l + CARD_PAD, top + Inches(0.05),
            col_w - 2 * CARD_PAD, Inches(0.35))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.name = FONT_HEADING
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
    row_h = (h - Inches(0.5)) / len(rows)
    for r_idx, row in enumerate(rows):
        ry = top + Inches(0.5) + r_idx * row_h
        for c_idx, cell_text in enumerate(row):
            l = MARGIN + c_idx * col_w
            bg_color = LIGHT_BG if r_idx % 2 == 0 else RGBColor(0xEC, 0xF0, 0xF5)
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                l, ry, col_w - Inches(0.05), row_h - Inches(0.03))
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg_color
            cell.line.fill.background()
            tb = slide.shapes.add_textbox(
                l + CARD_PAD, ry + Inches(0.05),
                col_w - 2 * CARD_PAD, row_h - Inches(0.1))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = cell_text
            p.font.name = FONT_BODY
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_TEXT
            p.alignment = PP_ALIGN.CENTER
    footer_bar(slide, prs, footer_text)
    return slide
